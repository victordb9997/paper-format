from __future__ import annotations

import base64
import io
import os
import tempfile
import uuid
from dataclasses import dataclass
from pathlib import Path
from typing import List

from flask import Flask, Response, flash, redirect, render_template, request, send_file, url_for
from werkzeug.utils import secure_filename

from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.lib.units import inch
from reportlab.lib.utils import ImageReader
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfgen import canvas


ALLOWED_DOCX = {".docx"}
ALLOWED_PPTX = {".pptx"}

app = Flask(__name__)
app.config["MAX_CONTENT_LENGTH"] = 150 * 1024 * 1024
app.secret_key = os.environ.get("SECRET_KEY", "dev")


@app.errorhandler(413)
def request_entity_too_large(error: Exception) -> Response:
    flash("Upload too large. Please compress your PPTX/DOCX or reduce image sizes and try again.")
    return redirect(url_for("index"))


@dataclass
class FigureAsset:
    filename: str
    caption: str
    data_uri: str


@dataclass
class FormattedPaper:
    paper_id: str
    title: str
    abstract: str
    body_sections: List[str]
    figures: List[FigureAsset]
    html: str
    pdf_bytes: bytes


STORE: dict[str, FormattedPaper] = {}


def _allowed_file(filename: str, allowed_exts: set[str]) -> bool:
    return Path(filename).suffix.lower() in allowed_exts


def _extract_docx(docx_path: Path) -> tuple[str, str, List[str]]:
    doc = Document(str(docx_path))
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    title = paragraphs[0] if paragraphs else "Untitled Manuscript"
    abstract = paragraphs[1] if len(paragraphs) > 1 else ""
    body_sections = paragraphs[2:] if len(paragraphs) > 2 else []
    return title, abstract, body_sections


def _extract_figures(pptx_path: Path) -> List[FigureAsset]:
    presentation = Presentation(str(pptx_path))
    figures: List[FigureAsset] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes, start=1):
            if not hasattr(shape, "image"):
                continue
            image = shape.image
            ext = image.ext
            filename = f"figure-{slide_index}-{shape_index}.{ext}"
            caption = f"Figure {slide_index}.{shape_index}"
            encoded = base64.b64encode(image.blob).decode("utf-8")
            data_uri = f"data:image/{ext};base64,{encoded}"
            figures.append(FigureAsset(filename=filename, caption=caption, data_uri=data_uri))
    return figures


def _render_html(title: str, abstract: str, body_sections: List[str], figures: List[FigureAsset]) -> str:
    return render_template(
        "preview.html",
        title=title,
        abstract=abstract,
        body_sections=body_sections,
        figures=figures,
        render_only=True,
    )


def _split_text(text: str, max_width: float, font_name: str, font_size: int) -> List[str]:
    words = text.split()
    if not words:
        return [""]
    lines: List[str] = []
    current = words[0]
    for word in words[1:]:
        test_line = f"{current} {word}"
        if pdfmetrics.stringWidth(test_line, font_name, font_size) <= max_width:
            current = test_line
        else:
            lines.append(current)
            current = word
    lines.append(current)
    return lines


def _draw_paragraph(
    pdf: canvas.Canvas,
    text: str,
    x: float,
    y: float,
    max_width: float,
    font_name: str,
    font_size: int,
    leading: int,
) -> float:
    pdf.setFont(font_name, font_size)
    for line in _split_text(text, max_width, font_name, font_size):
        if y < inch:
            pdf.showPage()
            y = letter[1] - inch
            pdf.setFont(font_name, font_size)
        pdf.drawString(x, y, line)
        y -= leading
    return y


def _create_pdf(title: str, abstract: str, body_sections: List[str], figures: List[FigureAsset]) -> bytes:
    buffer = io.BytesIO()
    pdf = canvas.Canvas(buffer, pagesize=letter)
    width, height = letter
    margin = inch
    current_y = height - margin

    pdf.setTitle(title or "Manuscript")

    current_y = _draw_paragraph(
        pdf,
        title or "Untitled Manuscript",
        margin,
        current_y,
        width - 2 * margin,
        "Helvetica-Bold",
        18,
        22,
    )
    current_y -= 10

    current_y = _draw_paragraph(
        pdf,
        "Abstract",
        margin,
        current_y,
        width - 2 * margin,
        "Helvetica-Bold",
        12,
        16,
    )
    current_y = _draw_paragraph(
        pdf,
        abstract,
        margin,
        current_y,
        width - 2 * margin,
        "Helvetica",
        11,
        15,
    )
    current_y -= 12

    for paragraph in body_sections:
        current_y = _draw_paragraph(
            pdf,
            paragraph,
            margin,
            current_y,
            width - 2 * margin,
            "Helvetica",
            11,
            15,
        )
        current_y -= 8

    if figures:
        current_y = _draw_paragraph(
            pdf,
            "Figures",
            margin,
            current_y,
            width - 2 * margin,
            "Helvetica-Bold",
            12,
            16,
        )
        for figure in figures:
            if current_y < 2 * inch:
                pdf.showPage()
                current_y = height - margin
            if figure.data_uri.startswith("data:image/"):
                header, encoded = figure.data_uri.split(",", 1)
                image_bytes = base64.b64decode(encoded)
                image = ImageReader(io.BytesIO(image_bytes))
                image_width, image_height = image.getSize()
                max_width = width - 2 * margin
                max_height = 3.25 * inch
                scale = min(max_width / image_width, max_height / image_height, 1.0)
                draw_width = image_width * scale
                draw_height = image_height * scale
                pdf.drawImage(
                    image,
                    margin,
                    current_y - draw_height,
                    width=draw_width,
                    height=draw_height,
                    preserveAspectRatio=True,
                    mask="auto",
                )
                current_y -= draw_height + 8
            current_y = _draw_paragraph(
                pdf,
                figure.caption,
                margin,
                current_y,
                width - 2 * margin,
                "Helvetica-Oblique",
                9,
                12,
            )
            current_y -= 10

    pdf.showPage()
    pdf.save()
    buffer.seek(0)
    return buffer.read()


@app.route("/", methods=["GET"])
def index() -> str:
    return render_template("index.html")


@app.route("/format", methods=["POST"])
def format_manuscript() -> Response:
    docx_file = request.files.get("manuscript")
    pptx_file = request.files.get("figures")

    if not docx_file or not pptx_file:
        flash("Please upload both a DOCX manuscript and a PPTX figures deck.")
        return redirect(url_for("index"))

    if not _allowed_file(docx_file.filename, ALLOWED_DOCX):
        flash("Manuscript must be a .docx file.")
        return redirect(url_for("index"))

    if not _allowed_file(pptx_file.filename, ALLOWED_PPTX):
        flash("Figures must be a .pptx file.")
        return redirect(url_for("index"))

    with tempfile.TemporaryDirectory() as tmpdir:
        docx_path = Path(tmpdir) / secure_filename(docx_file.filename)
        pptx_path = Path(tmpdir) / secure_filename(pptx_file.filename)
        docx_file.save(docx_path)
        pptx_file.save(pptx_path)

        title, abstract, body_sections = _extract_docx(docx_path)
        figures = _extract_figures(pptx_path)

        rendered_html = _render_html(title, abstract, body_sections, figures)
        pdf_bytes = _create_pdf(title, abstract, body_sections, figures)

    paper_id = uuid.uuid4().hex
    STORE[paper_id] = FormattedPaper(
        paper_id=paper_id,
        title=title,
        abstract=abstract,
        body_sections=body_sections,
        figures=figures,
        html=rendered_html,
        pdf_bytes=pdf_bytes,
    )

    return redirect(url_for("preview", paper_id=paper_id))


@app.route("/preview/<paper_id>", methods=["GET"])
def preview(paper_id: str) -> Response:
    paper = STORE.get(paper_id)
    if not paper:
        flash("We couldn't find that formatted paper. Please upload again.")
        return redirect(url_for("index"))
    return render_template(
        "preview.html",
        title=paper.title,
        abstract=paper.abstract,
        body_sections=paper.body_sections,
        figures=paper.figures,
        render_only=False,
        paper_id=paper.paper_id,
    )


@app.route("/download/<paper_id>", methods=["GET"])
def download(paper_id: str) -> Response:
    paper = STORE.get(paper_id)
    if not paper:
        flash("We couldn't find that formatted paper. Please upload again.")
        return redirect(url_for("index"))
    return send_file(
        io.BytesIO(paper.pdf_bytes),
        mimetype="application/pdf",
        as_attachment=True,
        download_name=f"{paper.title or 'manuscript'}.pdf",
    )


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
